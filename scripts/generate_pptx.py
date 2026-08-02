#!/usr/bin/env python3
"""Generate the AiCOD Solutions Limited company presentation as a .pptx file.

Usage:
    python3 scripts/generate_pptx.py

Requires:
    pip install python-pptx

Regenerate this file whenever the content in `presentation.html` changes so
the downloadable PowerPoint stays in sync with the on-site slide deck.
"""

import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
LOGO_PATH = os.path.join(ROOT, "AICOD SOLUTIONS LOGO.jpg")
OUTPUT_PATH = os.path.join(ROOT, "AiCOD_Solutions_Company_Presentation.pptx")

PRIMARY = RGBColor(0x00, 0x66, 0xFF)
SECONDARY = RGBColor(0xFF, 0x6B, 0x00)
DARK = RGBColor(0x1A, 0x1F, 0x2B)
LIGHT_GRAY = RGBColor(0xF4, 0xF7, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def add_background(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_gradient_bar(slide, top):
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), top, Inches(1.1), Pt(6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    bar.shadow.inherit = False


def add_title(slide, text, color=PRIMARY, top=Inches(0.5)):
    box = slide.shapes.add_textbox(Inches(0.6), top, Inches(12.1), Inches(1.0))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = color
    run.font.name = "Poppins"
    add_gradient_bar(slide, top + Inches(0.95))
    return box


def add_body_text(slide, text, top, left=Inches(0.6), width=Inches(12.1), size=20, color=DARK, bold=False):
    box = slide.shapes.add_textbox(left, top, width, Inches(1.5))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Plus Jakarta Sans"
    return box


def add_card(slide, left, top, width, height, title, body, icon_color=SECONDARY):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT_GRAY
    card.line.color.rgb = PRIMARY
    card.line.width = Pt(1.5)
    card.shadow.inherit = False

    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.margin_right = Inches(0.25)
    tf.margin_top = Inches(0.2)

    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = title
    r1.font.size = Pt(20)
    r1.font.bold = True
    r1.font.color.rgb = icon_color
    r1.font.name = "Poppins"

    p2 = tf.add_paragraph()
    p2.space_before = Pt(8)
    r2 = p2.add_run()
    r2.text = body
    r2.font.size = Pt(14)
    r2.font.color.rgb = DARK
    r2.font.name = "Plus Jakarta Sans"


def add_stat_box(slide, left, top, width, height, number, label):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.gradient()
    stops = box.fill.gradient_stops
    stops[0].color.rgb = PRIMARY
    stops[0].position = 0.0
    stops[1].color.rgb = SECONDARY
    stops[1].position = 1.0
    box.line.fill.background()
    box.shadow.inherit = False

    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = number
    r1.font.size = Pt(28)
    r1.font.bold = True
    r1.font.color.rgb = WHITE
    r1.font.name = "Poppins"

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = label
    r2.font.size = Pt(13)
    r2.font.color.rgb = WHITE
    r2.font.name = "Plus Jakarta Sans"


def add_bullet_list(slide, items, top, left=Inches(0.6), width=Inches(12.1)):
    box = slide.shapes.add_textbox(left, top, width, Inches(5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(14)
        run = p.add_run()
        run.text = f"\u2713  {item}"
        run.font.size = Pt(20)
        run.font.color.rgb = DARK
        run.font.name = "Plus Jakarta Sans"


def new_slide(prs, bg=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    add_background(slide, bg)
    return slide


def build_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Slide 1: Title
    slide = new_slide(prs, bg=WHITE)
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, Inches(5.72), Inches(1.0), height=Inches(1.6))
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.9), Inches(11.33), Inches(1.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "AiCOD Solutions Limited"
    r.font.size = Pt(44)
    r.font.bold = True
    r.font.color.rgb = DARK
    r.font.name = "Poppins"

    sub_box = slide.shapes.add_textbox(Inches(1), Inches(3.9), Inches(11.33), Inches(0.8))
    tf2 = sub_box.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "Empowering Businesses Through Innovative Digital Solutions"
    r2.font.size = Pt(20)
    r2.font.color.rgb = PRIMARY
    r2.font.name = "Plus Jakarta Sans"

    meta_box = slide.shapes.add_textbox(Inches(1), Inches(4.6), Inches(11.33), Inches(0.6))
    tf3 = meta_box.text_frame
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run()
    r3.text = "Company Overview & Presentation  •  Kabale, Uganda"
    r3.font.size = Pt(14)
    r3.font.color.rgb = SECONDARY
    r3.font.bold = True
    r3.font.name = "Plus Jakarta Sans"

    # Slide 2: Agenda
    slide = new_slide(prs)
    add_title(slide, "What We'll Cover")
    agenda = [
        "About AiCOD Solutions",
        "Mission, Vision & Values",
        "Services We Offer",
        "Strategic Objectives",
        "Organizational Structure",
        "Get In Touch",
    ]
    add_bullet_list(slide, agenda, Inches(1.9))

    # Slide 3: About Us
    slide = new_slide(prs)
    add_title(slide, "Who We Are")
    add_body_text(
        slide,
        "AiCOD Solutions Limited is a dynamic technology company based in Kabale, Uganda, "
        "specializing in comprehensive digital solutions for businesses and individuals. "
        "Founded on 1st January 2025, we help clients grow through modern digital techniques.",
        Inches(1.7),
        size=18,
    )
    stats = [("36+", "Businesses Helped / Year"), ("95%+", "Customer Satisfaction Target"), ("1000+", "Individuals Trained (2 yrs)")]
    left = Inches(0.6)
    box_width = Inches(3.9)
    gap = Inches(0.2)
    for i, (num, label) in enumerate(stats):
        add_stat_box(slide, Inches(0.6) + i * (box_width + gap), Inches(3.6), box_width, Inches(1.6), num, label)

    # Slide 4: Mission & Vision
    slide = new_slide(prs)
    add_title(slide, "Mission & Vision")
    add_card(
        slide, Inches(0.6), Inches(1.9), Inches(5.9), Inches(4.5),
        "Our Mission",
        "Empowering businesses and individuals through innovative digital solutions, "
        "exceptional service, and continuous learning, while fostering growth, "
        "productivity, and success in the digital age.",
    )
    add_card(
        slide, Inches(6.75), Inches(1.9), Inches(5.9), Inches(4.5),
        "Our Vision",
        "To be a leading technology company in Uganda, recognized for expertise in "
        "digital solutions, commitment to customer satisfaction, and contribution to "
        "the growth of the digital economy.",
    )

    # Slide 5: Core Values
    slide = new_slide(prs)
    add_title(slide, "Our Core Values")
    values = ["Innovation", "Customer-Centricity", "Excellence", "Collaboration", "Integrity"]
    left = Inches(0.6)
    top = Inches(2.2)
    for i, v in enumerate(values):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top + Inches(i * 0.9), Inches(6), Inches(0.7))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = PRIMARY
        box.line.width = Pt(2)
        box.shadow.inherit = False
        tf = box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = v
        r.font.size = Pt(18)
        r.font.bold = True
        r.font.color.rgb = PRIMARY
        r.font.name = "Plus Jakarta Sans"

    # Slide 6: Services
    slide = new_slide(prs)
    add_title(slide, "Services We Offer")
    services = [
        "Website Design",
        "Mobile App Development",
        "Advanced ICT Training",
        "Computer Repair & Maintenance",
        "Academic Research Assistance",
        "Security Camera Installation",
    ]
    cols, rows = 3, 2
    box_w, box_h = Inches(3.9), Inches(2.0)
    gap_x, gap_y = Inches(0.2), Inches(0.3)
    start_left, start_top = Inches(0.6), Inches(1.9)
    for i, s in enumerate(services):
        col = i % cols
        row = i // cols
        left = start_left + col * (box_w + gap_x)
        top = start_top + row * (box_h + gap_y)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, box_w, box_h)
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_GRAY
        box.line.fill.background()
        box.shadow.inherit = False
        tf = box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = s
        r.font.size = Pt(18)
        r.font.bold = True
        r.font.color.rgb = DARK
        r.font.name = "Plus Jakarta Sans"

    # Slide 7: Strategic Objectives
    slide = new_slide(prs)
    add_title(slide, "Strategic Objectives")
    objectives = [
        "Help 36 businesses adopt digital solutions every year",
        "Achieve a customer satisfaction rate of 95%+",
        "Launch 2 new digital products or services annually",
        "Grow the team by 20% within 12 months",
        "Provide digital skills training to 1,000 individuals in 2 years",
    ]
    add_bullet_list(slide, objectives, Inches(1.9))

    # Slide 8: Organizational Structure
    slide = new_slide(prs)
    add_title(slide, "Organizational Structure")
    ceo_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.67), Inches(1.9), Inches(4), Inches(0.8))
    ceo_box.fill.gradient()
    ceo_stops = ceo_box.fill.gradient_stops
    ceo_stops[0].color.rgb = PRIMARY
    ceo_stops[0].position = 0.0
    ceo_stops[1].color.rgb = SECONDARY
    ceo_stops[1].position = 1.0
    ceo_box.line.fill.background()
    ceo_box.shadow.inherit = False
    tf = ceo_box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "CEO / Managing Director"
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = "Plus Jakarta Sans"

    roles = [
        "CTO / Head of\nDevelopment",
        "Head of Sales\n& Marketing",
        "Head of Finance\n& Admin",
        "Head of Training\n& Research",
        "Head of IT Services\n& Support",
        "Head of Security\nSolutions",
    ]
    box_w, box_h = Inches(1.95), Inches(1.4)
    gap = Inches(0.15)
    total_w = 6 * box_w + 5 * gap
    start_left = (SLIDE_WIDTH - total_w) / 2
    top = Inches(3.4)
    for i, role in enumerate(roles):
        left = start_left + i * (box_w + gap)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, box_w, box_h)
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = PRIMARY
        box.line.width = Pt(1.5)
        box.shadow.inherit = False
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = role
        r.font.size = Pt(12)
        r.font.bold = True
        r.font.color.rgb = DARK
        r.font.name = "Plus Jakarta Sans"

    # Slide 9: Team
    slide = new_slide(prs)
    add_title(slide, "Our Team")
    add_card(slide, Inches(0.6), Inches(1.9), Inches(5.9), Inches(2.5), "N. Wyclife", "CEO / Web Developer")
    add_card(slide, Inches(6.75), Inches(1.9), Inches(5.9), Inches(2.5), "I. Alfred", "Full Stack Developer / Mobile App Development Expert")

    # Slide 10: Contact
    slide = new_slide(prs)
    add_title(slide, "Let's Work Together")
    contact_lines = [
        "Plot 22, Bushekwiire Road, Kabale (Mukikolegyi), Uganda",
        "+256-771097982, +256-782937175",
        "aicodsolutionsltd@gmail.com",
        "Mon-Sat: 9AM-6PM, Sun: 2PM-5PM",
    ]
    add_bullet_list(slide, contact_lines, Inches(1.9))

    prs.save(OUTPUT_PATH)
    print(f"Saved presentation to {OUTPUT_PATH}")


if __name__ == "__main__":
    build_presentation()
